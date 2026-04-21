"""
Rebuild the final/closing presentation with reviewer-recommended changes.
Preserves the visual design (decorative freeforms, fonts, colors) of the
original Overview.pptx by cloning existing slide templates and rewriting
only their text content.
"""
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = "ppt/Overview.pptx"
OUT = "ppt/Overview_Final.pptx"

TITLE_FONT = "Century Gothic Paneuropean Bold"
BODY_FONT = "Canva Sans"
ALT_BODY_FONT = "Century Gothic Paneuropean"


# ---------- helpers ----------
def clone_slide(prs, src_slide):
    """Duplicate a slide (including all shapes) and append to the end."""
    blank = prs.slide_layouts[6]  # Blank
    new_slide = prs.slides.add_slide(blank)
    # Copy shapes from src to new
    for shp in src_slide.shapes:
        el = shp.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return new_slide


def reorder_slides(prs, order):
    """Re-order the slides in the presentation according to index list."""
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    new_order = [slides[i] for i in order]
    for s in slides:
        sldIdLst.remove(s)
    for s in new_order:
        sldIdLst.append(s)


def set_title(slide, text, shape_name=None):
    """Set the title text on the first large/top text shape (or named shape)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape_name and shape.name != shape_name:
            continue
        # heuristic: top < 2in and font big
        if shape.top is not None and shape.top < Inches(2.3):
            tf = shape.text_frame
            # wipe current text but keep first paragraph formatting
            p0 = tf.paragraphs[0]
            # clear run texts, keep first run
            if p0.runs:
                first = p0.runs[0]
                first.text = text
                # remove other runs
                for r in p0.runs[1:]:
                    r._r.getparent().remove(r._r)
                # remove other paragraphs
                for extra in list(tf.paragraphs[1:]):
                    extra._p.getparent().remove(extra._p)
            else:
                p0.text = text
            return shape
    return None


def replace_body_text(slide, bullets, body_shape_idx=None, font_size=None):
    """Replace body bullets on the largest/lowest text shape that isn't the title."""
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.top and shape.top >= Inches(2.3):
            candidates.append(shape)
    if not candidates:
        return None
    # pick the largest (by area)
    if body_shape_idx is not None and body_shape_idx < len(candidates):
        body = candidates[body_shape_idx]
    else:
        body = max(candidates, key=lambda s: (s.width or 0) * (s.height or 0))

    tf = body.text_frame
    tf.word_wrap = True

    # capture first run formatting (template)
    tmpl_run = None
    for p in tf.paragraphs:
        if p.runs:
            tmpl_run = p.runs[0]
            break

    tmpl_font_name = tmpl_run.font.name if tmpl_run else BODY_FONT
    tmpl_font_size = tmpl_run.font.size if tmpl_run else Pt(28)
    if font_size is not None:
        tmpl_font_size = Pt(font_size)

    # Clear everything
    # remove all paragraphs
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)
    # create a fresh paragraph
    new_p = etree.SubElement(tf._txBody, "{http://schemas.openxmlformats.org/drawingml/2006/main}p")

    def add_para(text, bold=False, size=None, indent=0):
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        p = etree.SubElement(tf._txBody, f"{{{nsmap_a}}}p")
        pPr = etree.SubElement(p, f"{{{nsmap_a}}}pPr")
        pPr.set("lvl", str(indent))
        r = etree.SubElement(p, f"{{{nsmap_a}}}r")
        rPr = etree.SubElement(r, f"{{{nsmap_a}}}rPr")
        rPr.set("lang", "en-US")
        if bold:
            rPr.set("b", "1")
        sz = size if size is not None else tmpl_font_size
        if sz:
            rPr.set("sz", str(int(sz.pt * 100)) if hasattr(sz, 'pt') else str(int(sz * 100)))
        rPr.set("dirty", "0")
        latin = etree.SubElement(rPr, f"{{{nsmap_a}}}latin")
        latin.set("typeface", tmpl_font_name or BODY_FONT)
        solid = etree.SubElement(rPr, f"{{{nsmap_a}}}solidFill")
        srgb = etree.SubElement(solid, f"{{{nsmap_a}}}srgbClr")
        srgb.set("val", "000000")
        t = etree.SubElement(r, f"{{{nsmap_a}}}t")
        t.text = text
        return p

    # remove the placeholder empty paragraph
    tf._txBody.remove(new_p)

    for item in bullets:
        if isinstance(item, tuple):
            txt, opts = item
            add_para(txt, bold=opts.get("bold", False), size=opts.get("size"), indent=opts.get("indent", 0))
        else:
            add_para(item)
    return body


# ---------- build ----------
prs = Presentation(SRC)

# Original slides: 0..6
# 0: Title, 1: Problem, 2: Methodology, 3: Experiments, 4: MADOS, 5: Conclusion, 6: Thank You

title_tpl = prs.slides[0]
section_tpl = prs.slides[1]          # Problem background - big title + bullets
method_tpl = prs.slides[2]           # title + long bullet list
experiments_tpl = prs.slides[3]      # Two-block content
mados_tpl = prs.slides[4]            # Two-block content
conclusion_tpl = prs.slides[5]       # Two-block content
thankyou_tpl = prs.slides[6]         # closing

# --- We'll clone more slides from existing templates, then reorder at the end. ---

# Clone templates we'll need extras of.
# We need roughly 14 slides, we have 7 originals. Clone 7 more.
# We will use indexes below; clones are appended, so new indices start at 7.

# Clone plan (from -> new idx):
# 7  <- methodology tpl  -> Project Objectives
# 8  <- experiments tpl  -> Experiments Overview / Comparison
# 9  <- methodology tpl  -> Experiment 1 detail
# 10 <- methodology tpl  -> Experiment 2 detail
# 11 <- experiments tpl  -> Experiment 3 MADOS detail
# 12 <- experiments tpl  -> Post-processing & Sea Snot Index
# 13 <- experiments tpl  -> Field validation
# 14 <- methodology tpl  -> Key insights
# 15 <- experiments tpl  -> Future work

clones = []
clones.append(clone_slide(prs, method_tpl))       # 7 project objectives
clones.append(clone_slide(prs, experiments_tpl))  # 8 exp overview
clones.append(clone_slide(prs, method_tpl))       # 9 exp1 detail
clones.append(clone_slide(prs, method_tpl))       # 10 exp2 detail
clones.append(clone_slide(prs, experiments_tpl))  # 11 exp3 / MADOS detail
clones.append(clone_slide(prs, experiments_tpl))  # 12 SSI + post-processing
clones.append(clone_slide(prs, experiments_tpl))  # 13 field validation
clones.append(clone_slide(prs, method_tpl))       # 14 key insights
clones.append(clone_slide(prs, experiments_tpl))  # 15 future work

# Re-reference by index
slides = list(prs.slides)
S_TITLE       = slides[0]
S_PROBLEM     = slides[1]
S_METHOD      = slides[2]
S_EXP_OLD     = slides[3]       # will become unused / re-used as Exp Overview? We'll re-use.
S_MADOS_OLD   = slides[4]       # re-used for Exp 3 details
S_CONCL       = slides[5]
S_THANKS      = slides[6]
S_OBJ         = slides[7]
S_EXP_OVER    = slides[8]
S_EXP1        = slides[9]
S_EXP2        = slides[10]
S_EXP3        = slides[11]
S_POSTPROC    = slides[12]
S_FIELD       = slides[13]
S_KEY         = slides[14]
S_FUTURE      = slides[15]

# ---------- content ----------

# Slide 1: Title (minor enhancements)
# keep design but update under-guidance block to include focal persons line
# Find textbox 15 and 16 and tweak text
for shape in S_TITLE.shapes:
    if not shape.has_text_frame:
        continue
    if shape.name == "TextBox 3":  # JRF
        # leave as is
        pass
    if shape.name == "TextBox 15":  # Under Guidance
        tf = shape.text_frame
        # rewrite paragraphs
        for p in list(tf.paragraphs):
            p._p.getparent().remove(p._p)
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        def add(p_text, bold=False, size=28):
            p = etree.SubElement(tf._txBody, f"{{{nsmap_a}}}p")
            r = etree.SubElement(p, f"{{{nsmap_a}}}r")
            rPr = etree.SubElement(r, f"{{{nsmap_a}}}rPr")
            rPr.set("lang", "en-US")
            rPr.set("sz", str(int(size * 100)))
            if bold:
                rPr.set("b", "1")
            latin = etree.SubElement(rPr, f"{{{nsmap_a}}}latin")
            latin.set("typeface", "Canva Sans Bold" if bold else "Canva Sans")
            solid = etree.SubElement(rPr, f"{{{nsmap_a}}}solidFill")
            srgb = etree.SubElement(solid, f"{{{nsmap_a}}}srgbClr")
            srgb.set("val", "000000")
            t = etree.SubElement(r, f"{{{nsmap_a}}}t")
            t.text = p_text
        add("Under Guidance:", bold=True, size=30)
        add(" Dr. Madhuri Bhavsar (PI)", size=26)
        add(" Dr. Sudeep Tanwar, Dr. Swati Jain", size=26)
        add(" Dr. Jaiprakash Verma, Dr. Madhukant Patel", size=26)
        add("Focal Persons (SAC/ISRO):", bold=True, size=28)
        add(" Dr. Jishad M., Dr. Aditya Chaudhary,", size=24)
        add(" Dr. Surisetty V. V. Arunkumar, Dr. Neeraj Agrawal", size=24)

# Slide 2: Problem & Motivation - tighten + India context
set_title(S_PROBLEM, "PROBLEM BACKGROUND & MOTIVATION")
replace_body_text(
    S_PROBLEM,
    [
        ("Plastic makes up 60–80% of global marine litter — up to 95% in coastal hotspots.", {"size": 30}),
        ("India has a 7,500+ km coastline; rivers such as the Cooum and Ganga are major plastic outflows to the ocean.", {"size": 30}),
        ("Manual monitoring (boats, drones) is costly and spatially limited.", {"size": 30}),
        ("Satellites enable large-scale, low-cost, repeat observation of floating litter.", {"size": 30}),
        ("Goal: Build an AI/ML pipeline to detect floating plastic from multispectral remote sensing data.", {"bold": True, "size": 30}),
    ],
)

# Slide 3: Project Objectives (from methodology template)
set_title(S_OBJ, "PROJECT OBJECTIVES")
replace_body_text(
    S_OBJ,
    [
        ("1. Study floating marine debris using remote-sensing optical data.", {"bold": True, "size": 32}),
        ("2. Spectrally characterize plastic and build a spectral library.", {"bold": True, "size": 32}),
        ("3. Develop AI/ML models to automatically distinguish plastic from water, vegetation, sediment, etc.", {"bold": True, "size": 32}),
        ("4. Validate on satellite (Sentinel-2) imagery and real coastal sites.", {"bold": True, "size": 32}),
        ("", {"size": 20}),
        ("Target: Detect plastic even at ≥1% pixel coverage — the smallest realistic signal at satellite resolution.", {"bold": True, "size": 30}),
    ],
)

# Slide 4: Methodology pipeline
set_title(S_METHOD, "METHODOLOGY — END-TO-END PIPELINE")
replace_body_text(
    S_METHOD,
    [
        ("Data Sources", {"bold": True, "size": 30}),
        ("• Hyperspectral spectrometer (Exp 1) — 469 samples, 1%–100% pixel coverage.", {"size": 26}),
        ("• Reve multispectral sensor (Exp 2) — 18 bands, 410–940 nm, 79 → 400 samples (SMOTE).", {"size": 26}),
        ("• MADOS / Sentinel-2 satellite (Exp 3) — 2,803 images, 15 classes, ~1.5 M annotated pixels.", {"size": 26}),
        ("Processing Pipeline", {"bold": True, "size": 30}),
        ("• Preprocessing: reflectance normalization, ACOLITE atmospheric correction.", {"size": 26}),
        ("• Feature engineering: NDVI, FDI, proposed Sea Snot Index (SSI).", {"size": 26}),
        ("• Class balancing: SMOTE.", {"size": 26}),
        ("• Models: SVM, RF, MLP, kNN, AdaBoost, 1D-CNN, 1D-CNN + Transformer, Self-supervised.", {"size": 26}),
        ("• Post-processing: DBSCAN, Closing, Modulo Filter, SAM segmentation.", {"size": 26}),
    ],
)

# Slide 5: Experiments overview comparison
set_title(S_EXP_OVER, "EXPERIMENTS OVERVIEW")
replace_body_text(
    S_EXP_OVER,
    [
        ("Exp 1 — Hyperspectral Spectrometer", {"bold": True, "size": 30}),
        ("Data: 469 samples (262 plastic) · pixel coverage 1–100% · water / sediment / vegetation.", {"size": 24}),
        ("Best model: 1D CNN + Transformer — 97% overall · 100% at 1% pixel coverage.", {"size": 24}),
        ("Exp 2 — Reve Multispectral (410–940 nm, 18 bands)", {"bold": True, "size": 30}),
        ("Data: 79 → 400 samples after SMOTE · outdoor + indoor.", {"size": 24}),
        ("Best models: SVM & MLP — 100%; AdaBoost 99%; RF 98% (using NDVI + 940 nm).", {"size": 24}),
        ("Exp 3 — MADOS / Sentinel-2 (real satellite data)", {"bold": True, "size": 30}),
        ("Data: 2,803 images, 15 classes, ~4,696 Marine Debris pixels.", {"size": 24}),
        ("Best model: SVM — Recall 0.83 (pixel), 0.98 (image); Precision improved via post-processing.", {"size": 24}),
    ],
)

# Slide 6: Experiment 1 detail
set_title(S_EXP1, "EXPERIMENT 1 — HYPERSPECTRAL CHARACTERIZATION")
replace_body_text(
    S_EXP1,
    [
        ("Design", {"bold": True, "size": 28}),
        ("• Custom test-pad; 30° spectrometer FOV; 5.75-inch target radius.", {"size": 24}),
        ("• 469 spectra — 262 plastic (multi-color, transparent) + water, sediment, vegetation, soil, iron, concrete.", {"size": 24}),
        ("• Pixel coverage sweep: 1%, 25%, 50%, 75%, 100%.", {"size": 24}),
        ("Key Results", {"bold": True, "size": 28}),
        ("• 1D CNN + Transformer — 97% overall accuracy, 100% at 1% pixel coverage.", {"size": 24}),
        ("• Random Forest 90%; SVM (trained) 96%; raw SVM failed (0%).", {"size": 24}),
        ("• Spectral normalization removed time-of-day variation (morning vs. evening).", {"size": 24}),
        ("• SHAP highlighted a small set of wavelengths that drive plastic classification.", {"size": 24}),
        ("Key Insight", {"bold": True, "size": 28}),
        ("• Models trained on one sensor did not transfer to another — must train per target sensor.", {"size": 24}),
    ],
)

# Slide 7: Experiment 2 detail
set_title(S_EXP2, "EXPERIMENT 2 — REVE MULTISPECTRAL (410–940 nm)")
replace_body_text(
    S_EXP2,
    [
        ("Design", {"bold": True, "size": 28}),
        ("• 18-band Reve sensor; indoor + outdoor (Charodi Lake, Nirma campus).", {"size": 24}),
        ("• 9 object classes; 79 raw samples → 400 after SMOTE oversampling.", {"size": 24}),
        ("Feature Engineering", {"bold": True, "size": 28}),
        ("• NDVI (Landsat-8 analogue): (R860 − R645) / (R860 + R645).", {"size": 24}),
        ("• FDI adapted: (R810 + (R900 − R810)·0.56) − R860.", {"size": 24}),
        ("• Empirical study → NDVI + 940 nm separates plastic vs. non-plastic cleanly (2D scatter).", {"size": 24}),
        ("Results (All bands + NDVI)", {"bold": True, "size": 28}),
        ("• SVM & MLP — 100% · AdaBoost 99% · RF 98% · kNN/DT 97% · GNB 95%.", {"size": 24}),
        ("• Self-supervised (encoder-decoder + RF): 0.87 – 0.96 across seeds.", {"size": 24}),
    ],
)

# Slide 8: Experiment 3 - MADOS / Sentinel-2
set_title(S_EXP3, "EXPERIMENT 3 — MADOS / SENTINEL-2")
replace_body_text(
    S_EXP3,
    [
        ("Dataset & Scale", {"bold": True, "size": 30}),
        ("• 2,803 Sentinel-2 images (240×240), 15 classes, ~1.5 M annotated pixels.", {"size": 26}),
        ("• 4,696 Marine Debris pixels; 7.6 L train / 3.5 L val / 3.6 L test pixels.", {"size": 26}),
        ("Pixel-level Results (SVM, best of 6 models)", {"bold": True, "size": 30}),
        ("• Recall 0.83 · Precision 0.60 · F1 0.70 (MD class).", {"size": 26}),
        ("Image-level Results (SVM, full dataset)", {"bold": True, "size": 30}),
        ("• TP 121 · FP 194 · TN 405 · FN 2 · Recall 0.98 · Precision 0.38 · F1 0.55.", {"size": 26}),
        ("• High recall, but false positives from ships / islands / Sea Snot drag precision.", {"size": 26}),
    ],
)

# Slide 9: Post-processing & SSI
set_title(S_POSTPROC, "IMPROVING PRECISION — SSI & POST-PROCESSING")
replace_body_text(
    S_POSTPROC,
    [
        ("Sea Snot Index (proposed)", {"bold": True, "size": 30}),
        ("• SSI = (R_B8 − R_B11) / (R_B8 + R_B11) — targets the spectral gap of sea snot.", {"size": 26}),
        ("• Motivation: sea snot was the largest source of false positives.", {"size": 26}),
        ("Post-processing — Image-level F1 on MADOS", {"bold": True, "size": 30}),
        ("• SVM baseline         Recall 0.98 · Precision 0.38 · F1 0.55", {"size": 26}),
        ("• + Closing (3×3)      Recall 0.05 · Precision 0.31 · F1 0.09", {"size": 26}),
        ("• + Modulo Filter (2×2) Recall 0.48 · Precision 0.49 · F1 0.48", {"size": 26}),
        ("• + SAM (Meta)         Recall 0.69 · Precision 0.42 · F1 0.52", {"size": 26}),
        ("• Takeaway: SAM best balances recall & precision; closing over-erodes small debris.", {"bold": True, "size": 26}),
    ],
)

# Slide 10: Field validation
set_title(S_FIELD, "FIELD VALIDATION — REAL SENTINEL-2 SITES")
replace_body_text(
    S_FIELD,
    [
        ("Sites Tested", {"bold": True, "size": 30}),
        ("• Limassol Bay (Cyprus) — plastic-net target (Jun 2025).", {"size": 26}),
        ("• Cooum River, Chennai — urban plastic outflow (Jun 2022).", {"size": 26}),
        ("• Ocean CleanUp Interceptors — river-mouth debris capture sites (Jan 2025).", {"size": 26}),
        ("Observations", {"bold": True, "size": 30}),
        ("• Model correctly localises debris patches; morphology preserved after SAM.", {"size": 26}),
        ("• Initial accuracy poor without atmospheric correction.", {"size": 26}),
        ("• After ACOLITE correction: predictions cleaner, noise drastically reduced.", {"size": 26}),
        ("• ACOLITE integration currently blocked by a ‘NoneType’ pre-processor error — in progress.", {"size": 26}),
    ],
)

# Slide 11: Key insights
set_title(S_KEY, "KEY INSIGHTS")
replace_body_text(
    S_KEY,
    [
        ("1. Plastic is detectable at ≥1% pixel coverage — 1D CNN + Transformer generalises across coverage levels.", {"bold": True, "size": 28}),
        ("2. NDVI + 940 nm is a surprisingly strong 2-feature separator at multispectral resolution.", {"bold": True, "size": 28}),
        ("3. Sea snot is the dominant false-positive class; a dedicated SSI helps but is not a full fix.", {"bold": True, "size": 28}),
        ("4. Post-processing trades recall for precision — SAM segmentation is the best compromise.", {"bold": True, "size": 28}),
        ("5. Atmospheric correction (ACOLITE) is non-negotiable for real Sentinel-2 scenes.", {"bold": True, "size": 28}),
        ("6. Cross-sensor transfer fails — models must be trained on data from the target sensor.", {"bold": True, "size": 28}),
    ],
)

# Slide 12: Conclusion
set_title(S_CONCL, "CONCLUSION")
replace_body_text(
    S_CONCL,
    [
        ("• Delivered a full pipeline: spectrum → features → ML/AI → post-processing → satellite deployment.", {"size": 28}),
        ("• Three complementary experiments covering sensor, platform and scale.", {"size": 28}),
        ("• 97–100% accuracy on controlled spectra; 0.98 recall on real MADOS imagery.", {"size": 28}),
        ("• Proposed a novel Sea Snot Index to attack the main false-positive source.", {"size": 28}),
        ("• Validated end-to-end on three real coastal sites with Sentinel-2 data.", {"size": 28}),
        ("", {"size": 14}),
        ("Outcome: A deployable, India-ready framework for satellite-based plastic monitoring.", {"bold": True, "size": 30}),
    ],
)

# Slide 13: Future work
set_title(S_FUTURE, "FUTURE WORK")
replace_body_text(
    S_FUTURE,
    [
        ("Short term", {"bold": True, "size": 30}),
        ("• Resolve ACOLITE pipeline error; integrate as an automated preprocessing step.", {"size": 26}),
        ("• Retrain detection models on Sentinel-2 bands to close the cross-sensor gap.", {"size": 26}),
        ("Medium term", {"bold": True, "size": 30}),
        ("• Release an open hyperspectral library of Indian plastics + environments.", {"size": 26}),
        ("• Extend the Sea Snot Index and test deep segmentation models (U-Net, DeepLabv3+).", {"size": 26}),
        ("Long term", {"bold": True, "size": 30}),
        ("• Operational deployment along the Indian coast; integrate with ISRO Sentinel-2 products.", {"size": 26}),
        ("", {"size": 14}),
        ("Ask: access to SAC / ISRO Sentinel-2 archive for target sites along the Indian coast.", {"bold": True, "size": 28}),
    ],
)

# Slide 14: Thank You - keep, maybe add contact line
# Find title and ensure it says THANK YOU
set_title(S_THANKS, "THANK YOU")

# ----- clean up the unused original Experiments / MADOS slides -----
# They have been repurposed:
#   slides[3] (S_EXP_OLD)  -> will be used as the Experiments Overview? We already used S_EXP_OVER (cloned 8).
#   slides[4] (S_MADOS_OLD) -> we already populated MADOS details on cloned S_EXP3.
# So we should DROP slides[3] and slides[4] from the final order.
# Keep them out of order instead.

# ---------- final order ----------
# Original indexes:
# 0 title, 1 problem, 2 methodology, 3 exp(old), 4 mados(old),
# 5 conclusion, 6 thankyou, 7 obj, 8 exp_over, 9 exp1, 10 exp2,
# 11 exp3, 12 postproc, 13 field, 14 key, 15 future
final_order = [
    0,  # title
    1,  # problem
    7,  # objectives
    2,  # methodology pipeline
    8,  # experiments overview
    9,  # exp 1
    10, # exp 2
    11, # exp 3 MADOS
    12, # SSI + post-processing
    13, # field validation
    14, # key insights
    5,  # conclusion
    15, # future work
    6,  # thank you
]
# The leftover slides (3, 4) will be deleted.

total = len(list(prs.slides))
# Re-order first (keeping all slides referenced)
reorder_slides(prs, final_order + [i for i in range(total) if i not in final_order])

# Now delete slides not in final_order by name (they're now at the tail)
# Simpler: remove by xml id of sldId elements not in final_order
sldIdLst = prs.slides._sldIdLst
keep_count = len(final_order)
all_sldIds = list(sldIdLst)
for s in all_sldIds[keep_count:]:
    rId = s.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    # drop the relationship and the element
    prs.part.drop_rel(rId)
    sldIdLst.remove(s)

prs.save(OUT)
print(f"Saved: {OUT} — {len(list(prs.slides))} slides")
